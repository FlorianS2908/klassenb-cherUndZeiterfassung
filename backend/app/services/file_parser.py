from __future__ import annotations

import html
import re
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from PyPDF2 import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".ppt", ".html", ".htm", ".md", ".txt"}


def detect_file_type(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return "PDF", "Seiten"
    if suffix in {".pptx", ".ppt"}:
        return "PowerPoint", "Folien"
    if suffix in {".html", ".htm"}:
        return "HTML", "Abschnitte"
    if suffix == ".md":
        return "Markdown", "Abschnitte"
    if suffix == ".txt":
        return "Text", "Zeilen"
    raise ValueError("Dateityp wird nicht unterstuetzt.")


def count_items(path: Path) -> int:
    file_type, _ = detect_file_type(path)
    if file_type == "PDF":
        return len(PdfReader(str(path)).pages)
    if file_type == "PowerPoint":
        if path.suffix.lower() == ".ppt":
            return 1
        return len(Presentation(str(path)).slides)
    if file_type in {"HTML", "Markdown"}:
        return len(split_sections(path)) or 1
    return len(path.read_text(encoding="utf-8", errors="ignore").splitlines()) or 1


def split_sections(path: Path) -> list[str]:
    text = path.read_text(encoding="utf-8", errors="ignore")
    if path.suffix.lower() in {".html", ".htm"}:
        soup = BeautifulSoup(text, "html.parser")
        headings = soup.find_all(re.compile("^h[1-6]$"))
        if not headings:
            return [soup.get_text("\n", strip=True)]
        sections: list[str] = []
        for heading in headings:
            chunk = [heading.get_text(" ", strip=True)]
            for sibling in heading.find_next_siblings():
                if sibling.name and re.fullmatch(r"h[1-6]", sibling.name):
                    break
                chunk.append(sibling.get_text(" ", strip=True) if hasattr(sibling, "get_text") else str(sibling))
            sections.append("\n".join(part for part in chunk if part))
        return sections
    parts = re.split(r"(?m)^#{1,6}\s+", text)
    headings = re.findall(r"(?m)^#{1,6}\s+(.+)$", text)
    if len(parts) <= 1:
        return [text]
    sections = []
    for idx, part in enumerate(parts[1:]):
        title = headings[idx] if idx < len(headings) else f"Abschnitt {idx + 1}"
        sections.append(f"{title}\n{part}".strip())
    return sections


def extract_text(path: Path, selected_items: list[int]) -> str:
    file_type, _ = detect_file_type(path)
    if file_type == "PDF":
        reader = PdfReader(str(path))
        return "\n\n".join((reader.pages[i - 1].extract_text() or "") for i in selected_items)
    if file_type == "PowerPoint":
        if path.suffix.lower() == ".ppt":
            return "Legacy-PPT-Dateien koennen nur eingeschraenkt ausgelesen werden. Bitte als PPTX speichern."
        slides = Presentation(str(path)).slides
        texts = []
        for item in selected_items:
            slide = slides[item - 1]
            texts.append("\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text")))
        return "\n\n".join(texts)
    if file_type in {"HTML", "Markdown"}:
        sections = split_sections(path)
        return "\n\n".join(html.unescape(sections[i - 1]) for i in selected_items)
    lines = path.read_text(encoding="utf-8", errors="ignore").splitlines()
    return "\n".join(lines[i - 1] for i in selected_items if i - 1 < len(lines))
